#!/usr/bin/env python3
# VeroRun 维洛智能 (verorun.com / verorun.cn)
# 版权所有 (c) 2026 樊聚民 (fanjumin). All Rights Reserved.

"""
Agent Matrix — Flask Blueprint
============================
注册为 /admin/agent-matrix/ 前缀。

API 端点统计: ~35 个
"""
import os, sys, json, logging

from i18n import _
from flask import Blueprint, request, jsonify, send_from_directory

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(BASE_DIR, '..', 'auth-center'))
sys.path.insert(0, os.path.join(BASE_DIR, '..'))

from services.jwt_service import validate_token
from models import get_db

agent_matrix_bp = Blueprint('agent_matrix', __name__, url_prefix='/admin/agent-matrix')

# 延迟加载 models（循环依赖处理）
_models = None
def _m():
    global _models
    if _models is None:
        from agent_matrix import models as m
        _models = m
    return _models


# ============================================================
# RAG 知识检索（Read 层）
# ============================================================

def _inject_knowledge(user_message: str, top_k: int = 5, scope: str = None) -> str:
    """
    基于用户输入实时检索 knowledge_blocks，拼入 system prompt。
    已升级为混合检索（向量+关键词+RRF），实现委托 agent_matrix/rag_retriever.py。
    返回格式化的知识文本，无结果返回空字符串。
    scope: 可选 'system' | 'user'，默认 None（检索全部）
    """
    try:
        from agent_matrix.rag_retriever import rag_search
        top = rag_search(user_message, top_k=top_k, scope=scope)

        if not top:
            return ''

        # 更新命中计数
        with _m().get_db() as conn:
            for block in top:
                if block.get('score', 0) > 0.3:
                    try:
                        conn.execute(
                            "UPDATE knowledge_blocks SET hit_count = hit_count + 1 WHERE id = %s",
                            (block['id'],)
                        )
                    except Exception:
                        pass
            conn.commit()

        # 拼成文本
        lines = []
        for block in top:
            if block.get('score', 0) > 0.3:
                lines.append(
                    f"- [{block['category']}] {block['title']}: {block['content'][:200]}"
                )

        if not lines:
            return ''

        return '\n\n=== 知识库（自动检索） ===\n' + '\n'.join(lines) + '\n=== 知识库结束 ==='

    except Exception as e:
        import logging
        logging.getLogger(__name__).warning(f"RAG 检索失败: {e}")
        return ''


# ============================================================
# 鉴权
# ============================================================

def _require_admin():
    auth = request.headers.get('Authorization', '')
    token = auth.replace('Bearer ', '') if auth.startswith('Bearer ') else auth
    if not token:
        token = request.cookies.get('sso_token') or request.cookies.get('tm_token')
    payload = validate_token(token) if token else None
    if not payload or not payload.get('is_admin'):
        return None, (jsonify({'success': False, 'error': _('Requires management permissions')}), 401)
    return payload, None


def _success(data=None, message='ok'):
    return jsonify({'success': True, 'data': data, 'message': message})


def _error(message, code=400):
    return jsonify({'success': False, 'error': message}), code


def _check_ai_access():
    """
    检查 AI 功能是否可用（独立部署订阅过期检查）
    仅在客户端模式（APP_MODE=client）生效
    返回 None 表示可用，返回 Response 表示已过期
    P1-F02: 改为 fail-closed，无法确认授权状态时返回 503
    """
    if os.environ.get('APP_MODE', 'main') != 'client':
        return None
    try:
        from services.license_service import LicenseService
        ls = LicenseService()
        if not ls.check_ai_access():
            return jsonify({
                'success': False,
                'error': _('Subscription expired, AI features unavailable'),
                'code': 'subscription_expired',
                'action': 'renew',
            }), 403
    except ImportError:
        import logging
        logging.getLogger(__name__).error(
            "[AgentMatrix] LicenseService import failed, blocking AI access (fail-closed)"
        )
        return jsonify({
            'success': False,
            'error': _('License service unavailable, AI features temporarily disabled'),
            'code': 'license_unavailable',
        }), 503
    except Exception as e:
        import logging
        logging.getLogger(__name__).error(
            f"[AgentMatrix] _check_ai_access error (fail-closed): {e}"
        )
        return jsonify({
            'success': False,
            'error': _('AI service unavailable, please try again later'),
            'code': 'service_error',
        }), 503
    return None


# ============================================================
# 1. Agent 管理
# ============================================================

@agent_matrix_bp.route('/agents', methods=['GET'])
def list_agents():
    admin, err = _require_admin()
    if err: return err

    role_type = request.args.get('role')
    domain = request.args.get('domain')
    active_only = request.args.get('active_only', '').lower() == 'true'

    agents = _m().list_agents(role_type=role_type, domain=domain, active_only=active_only)
    return _success(agents)


@agent_matrix_bp.route('/agents', methods=['POST'])
def create_agent():
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    if not data.get('name'):
        return _error(_('Agent name cannot be empty'))

    try:
        provider = data.get('provider', 'dashscope')
        # 处理 API Key 值——仅加密存入 provider_api_keys
        if 'api_key' in data:
            api_key_val = data.get('api_key', '')
            # 同步写入 provider_api_keys 表
            if api_key_val and provider:
                try:
                    from services.crypto import encrypt as _enc
                    encrypted = _enc(api_key_val)
                    with get_db() as conn:
                        row = conn.execute(
                            "INSERT INTO provider_api_keys (name, key_value_enc, provider, description) "
                            "VALUES (%s, %s, %s, %s) "
                            "ON CONFLICT (name, provider) DO UPDATE SET key_value_enc=EXCLUDED.key_value_enc, updated_at=NOW() "
                            "RETURNING id",
                            (f'agent:{data["name"]}', encrypted, provider, f'Agent: {data["name"]}')
                        ).fetchone()
                        if row:
                            data['api_key_id'] = row['id']
                        conn.commit()
                except Exception:
                    pass  # 写入 provider_api_keys 失败不影响 Agent 创建
            if 'api_key_ref' not in data:
                data['api_key_ref'] = ''
            del data['api_key']
        agent_id = _m().create_agent(data)
        return _success({'id': agent_id}, _('Agent has been created'))
    except Exception as e:
        return _error(str(e), 500)


@agent_matrix_bp.route('/agents/<int:aid>', methods=['GET'])
def get_agent(aid):
    admin, err = _require_admin()
    if err: return err

    agent = _m().get_agent(aid)
    if not agent:
        return _error(_('Agent does not exist'), 404)
    return _success(agent)


@agent_matrix_bp.route('/agents/<int:aid>', methods=['PUT'])
def update_agent(aid):
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    # 处理 API Key 值——仅加密存入 provider_api_keys
    if 'api_key' in data:
        agent = _m().get_agent(aid)
        provider = data.get('provider') or (agent.get('provider', 'dashscope') if agent else 'dashscope')
        api_key_val = data.get('api_key', '')
        # 同步写入 provider_api_keys 表
        if api_key_val and provider:
            try:
                from services.crypto import encrypt as _enc
                encrypted = _enc(api_key_val)
                agent_name = data.get('name') or (agent.get('name', '') if agent else '')
                name = f'agent:{agent_name}' if agent_name else f'{provider}_auto_key'
                with get_db() as conn:
                    row = conn.execute(
                        "INSERT INTO provider_api_keys (name, key_value_enc, provider, description) "
                        "VALUES (%s, %s, %s, %s) "
                        "ON CONFLICT (name, provider) DO UPDATE SET key_value_enc=EXCLUDED.key_value_enc, updated_at=NOW() "
                        "RETURNING id",
                        (name, encrypted, provider, f'Agent: {agent_name}')
                    ).fetchone()
                    if row:
                        data['api_key_id'] = row['id']
                    conn.commit()
            except Exception:
                pass  # 写入 provider_api_keys 失败不影响 Agent 更新
        if 'api_key_ref' not in data:
            data['api_key_ref'] = ''
        del data['api_key']
    ok = _m().update_agent(aid, data)
    if not ok:
        return _error(_('No fields to update'))
    return _success(None, _('Agent has been updated'))


@agent_matrix_bp.route('/agents/<int:aid>', methods=['DELETE'])
def delete_agent(aid):
    admin, err = _require_admin()
    if err: return err

    agent = _m().get_agent(aid)
    if not agent:
        return _error(_('Agent does not exist'), 404)
    if agent['role_type'] == 'master':
        return _error(_('Cannot delete main Agent'), 400)

    _m().delete_agent(aid)
    return _success(None, _('Agent has been deleted'))


@agent_matrix_bp.route('/agents/<int:aid>/toggle', methods=['POST'])
def toggle_agent(aid):
    admin, err = _require_admin()
    if err: return err

    new_state = _m().toggle_agent(aid)
    if new_state is None:
        return _error(_('Agent does not exist'), 404)
    return _success({'is_active': new_state}, f'Agent has been {"enabled" if new_state else "disabled"}')


@agent_matrix_bp.route('/agents/<int:aid>/test', methods=['POST'])
def test_agent(aid):
    """测试 Agent：发一条消息给它"""
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    query = data.get('query', '')
    if not query:
        return _error('请先输入测试消息（不能为空）')

    agent_config = _m().get_agent(aid)
    if not agent_config:
        return _error(_('Agent does not exist'), 404)

    from agent_matrix.agent_runner import AgentRunner
    runner = AgentRunner(agent_config)

    result = runner.execute({
        'task_id': 'test-' + str(aid),
        'title': _('Test Task'),
        'description': query,
        'input_data': {'query': query},
        'expected_output': {'fields': ['response']},
        'max_retries': 0,
    })

    return _success({
        'response': result.get('response', ''),
        'confidence': result.get('confidence', 0),
        'status': result.get('status', 'failed'),
        'logs': result.get('logs', [])
    })


@agent_matrix_bp.route('/agents/<int:aid>/capabilities', methods=['GET'])
def agent_capabilities(aid):
    admin, err = _require_admin()
    if err: return err

    agent = _m().get_agent(aid)
    if not agent:
        return _error(_('Agent does not exist'), 404)

    caps = {
        'name': agent['name'],
        'role_type': agent['role_type'],
        'domain': agent['domain'],
        'managed_modules': agent.get('managed_modules', '[]'),
        'capabilities': agent.get('capabilities', '[]'),
        'provider': agent['provider'],
        'model': agent['model_name'],
        'is_active': agent['is_active'],
    }
    return _success(caps)


# ============================================================
# 2. 任务管理
# ============================================================

@agent_matrix_bp.route('/tasks', methods=['GET'])
def list_tasks():
    admin, err = _require_admin()
    if err: return err

    status = request.args.get('status')
    module = request.args.get('module')
    agent_id = request.args.get('agent_id', type=int)
    limit = request.args.get('limit', 50, type=int)

    tasks = _m().list_tasks(status=status, module=module, agent_id=agent_id, limit=limit)
    return _success(tasks)


@agent_matrix_bp.route('/tasks/<task_id>', methods=['GET'])
def get_task(task_id):
    admin, err = _require_admin()
    if err: return err

    task = _m().get_task(task_id)
    if not task:
        return _error(_('Task does not exist'), 404)

    # 附带子任务
    task['sub_tasks'] = _m().get_sub_tasks(task_id)
    task['logs'] = _m().get_task_logs(task_id)
    return _success(task)


@agent_matrix_bp.route('/tasks/<task_id>/cancel', methods=['POST'])
def cancel_task(task_id):
    admin, err = _require_admin()
    if err: return err

    _m().cancel_task(task_id)
    return _success(None, _('Task has been canceled'))


@agent_matrix_bp.route('/tasks/<task_id>/retry', methods=['POST'])
def retry_task(task_id):
    """重试失败的任务"""
    admin, err = _require_admin()
    if err: return err

    task = _m().get_task(task_id)
    if not task:
        return _error(_('Task does not exist'), 404)

    _m().update_task_status(task_id, 'pending', error_message='')
    return _success(None, _('Task has been reset to pending'))


@agent_matrix_bp.route('/tasks/<task_id>/logs', methods=['GET'])
def task_logs(task_id):
    admin, err = _require_admin()
    if err: return err

    logs = _m().get_task_logs(task_id)
    return _success(logs)


@agent_matrix_bp.route('/tasks/recent', methods=['GET'])
def recent_tasks():
    admin, err = _require_admin()
    if err: return err

    tasks = _m().get_recent_tasks(limit=20)
    return _success(tasks)


# ============================================================
# 3. Master Agent 对话
# ============================================================

@agent_matrix_bp.route('/chat', methods=['POST'])
def chat_with_master():
    """
    向 Master Agent 发送指令（核心入口）

    请求: {
        "message": _("Write an intelligent article for me..."),
        "session_id": _("SESSION-... (Optional, leave blank to create a new session)")
    }
    """
    admin, err = _require_admin()
    if err: return err

    ai_err = _check_ai_access()
    if ai_err: return ai_err

    data = request.get_json(force=True) or {}
    message = data.get('message', '').strip()
    if not message:
        return _error(_('Message cannot be empty'))

    session_id = data.get('session_id', '')
    if not session_id:
        session_id = _m().create_session()
    # P1-F09: session_id 白名单校验（防路径穿越）
    elif not _m().is_valid_session_id(session_id):
        return _error(_('Invalid session ID'), 400)

    mode = data.get('mode', 'fast')

    # 获取 Master Agent
    agents = _m().list_agents(role_type='master', active_only=True)
    if not agents:
        return _error('没有可用的 Master Agent，请先创建', 500)

    master = agents[0]
    master_id = master['id']

    # 创建 Orchestrator 并执行
    from agent_matrix.orchestrator import AgentOrchestrator
    orchestrator = AgentOrchestrator(models_module=_m())

    try:
        result = orchestrator.process_instruction(
            instruction=message,
            master_agent_id=master_id,
            session_id=session_id,
            mode=mode,
            user_id=admin.get('user_id', 0)
        )
        # Write 层：对话结束自动提取知识（异步，不阻塞响应）
        orchestrator._on_task_complete(message, admin.get('user_id', 0), result)
    except Exception as e:
        import traceback
        traceback.print_exc()
        return _error(f'Execution failed: {e}', 500)

    return _success({
        'session_id': session_id,
        'master_task_id': result.get('master_task_id', ''),
        'summary': result.get('summary', ''),
        'decomposition': result.get('decomposition', []),
        'sub_task_results': result.get('sub_task_results', []),
        'status': result.get('status', 'failed'),
        'duration_s': result.get('duration_s', 0),
        'error': result.get('error', ''),  # BUG-001 诊断：暴露异常信息
    })


@agent_matrix_bp.route('/chat/tool', methods=['POST'])
def chat_tool():
    """工具调用模式：AI 分析意图后直接调用后端功能（PPT/图像/多媒体）"""
    admin, err = _require_admin()
    if err: return err

    ai_err = _check_ai_access()
    if ai_err: return ai_err

    data = request.get_json(force=True) or {}
    message = data.get('message', '').strip()
    if not message:
        return _error(_('Message cannot be empty'))

    session_id = data.get('session_id', '')
    if not session_id:
        session_id = _m().create_session()
    # P1-F09: session_id 白名单校验（防路径穿越）
    elif not _m().is_valid_session_id(session_id):
        return _error(_('Invalid session ID'), 400)

    # 获取 Master Agent 用于 AI 分析意图
    agents = _m().list_agents(role_type='master', active_only=True)
    if not agents:
        return _error(_('No available Master Agent'), 500)
    master = agents[0]

    # 工具路由映射
    intent_prompt = """你是一个意图分析器。分析用户消息，只返回一个 JSON：
{"intent":"ppt|image|voice|video|cms|supply_chain|clean|site_build|ads|chat","args":{}}

- ppt: 生成PPT。args: {topic, pages(默认10), style(可选)}
- image: 生成/分析图像。args: {prompt, style(可选), count(默认1), action:"generate"|"analyze"}
- voice: 克隆声音。args: {name, audio_url}
- video: 生成数字人视频。args: {title, text, voice_id(若有), image_url(可选)}
- cms: 写文章。args: {title, category(可选), content_prompt}
- supply_chain: 供应链与商城操作。args: {action:"search"|"collect"|"optimize"|"publish",keywords(可选),item_id(可选)}
- clean: 数据清洗。用户提供了需要清洗的原始内容（文章、白皮书、行业背景等）。args: {content: 原始内容全文}
- site_build: 用户想创建一个全新的网站（不是广告/文章/PPT等具体内容）。关键词包括_("Build a website")_("Create Website")_("Build Website")_("Help me build a website")_("Generate Website")。注意：创建广告/创建文章/生成PPT等具体内容操作不属于site_build。args: {prompt_identifier: 行业标识(如law_firm/restaurant等，从用户描述推断), action:"preview"|"execute"|"modify"}
- ads: 广告管理操作（优先级高于site_build）。用户提到了广告相关内容：_("Advertisement")_("Ad Position")"AD""banner"_("Campaign")_("Analyze Ad")_("Add Advertisement")"创建广告""新增广告""广告管理""查看广告""广告列表""帮我创建广告""生成广告""广告代码""ad_code"。注意：如果用户说"创建""新增""帮我创建""生成"广告，action必须设为create；如果说"查看""列出""查询"广告，action设为list。args: {action:"list"|"create"|"update"|"delete"|"stats"|"analyze"|"snippet", name, position, ad_type, image_url, link_url, ad_code, site_key(默认default), page(默认*), ad_id, days(默认7)}
- chat: 普通对话，不是工具调用。"""

    # 用轻量模型快速识别意图
    import json, re
    try:
        from agent_matrix.engine import UnifiedLLM
        engine = UnifiedLLM(master)
        intent_raw = engine.chat([
            {"role": "system", "content": intent_prompt},
            {"role": "user", "content": message}
        ], temperature=0.1, max_tokens=200)
        match = re.search(r'\{[\s\S]*\}', intent_raw)
        intent_data = json.loads(match.group(0)) if match else {}
        intent = intent_data.get('intent', 'chat')
        args = intent_data.get('args', {})
    except Exception:
        intent = 'chat'
        args = {}

    actions = []
    summary = ''

    try:
        if intent == 'ppt':
            topic = args.get('topic', message)
            pages = int(args.get('pages', 10))
            style = args.get('style', _('Dark Tech Style, 16:9'))
            filename = _generate_ppt_file(topic, pages, style)
            if filename:
                url = f"/admin/agent-matrix/media/download/{filename}"
                actions.append({'type': 'ppt_download', 'url': url, 'filename': filename})
                summary = f'✅ PPT已生成：{topic}（{pages}页）\n⬇ 点击下方按钮下载'
            else:
                summary = f'❌ PPT生成失败，请检查后端日志'

        elif intent == 'image':
            action_type = args.get('action', 'generate')
            if action_type == 'analyze':
                summary = '🔍 图像理解：请先上传图片（点击📎按钮），然后重新发送指令。'
            else:
                prompt = args.get('prompt', message)
                style = args.get('style', 'realistic')
                count = int(args.get('count', 1))
                size = args.get('size', '1024x1024')
                from agent_matrix.tools import execute_tool
                result = execute_tool('generate_image', {
                    'prompt': prompt, 'style': style,
                    'count': count, 'size': size
                })
                # 解析结果，提取图片 URL
                summary = result
                lines = result.split('\n')
                for ll in lines:
                    ll = ll.strip()
                    if ll.startswith('http') or ll.startswith('/admin'):
                        actions.append({'type': 'image', 'url': ll})

        elif intent == 'voice':
            name = args.get('name', _('My Voice'))
            audio_url = args.get('audio_url', '')
            if audio_url:
                # voice_clone 功能暂未实现，返回提示
                summary = f'🎙️ 声音克隆：{name}\n语音克隆功能开发中，请到「🎙️ 多媒体」Tab使用音频上传功能。'
                logger = logging.getLogger(__name__)
                logger.warning(f"[Voice] voice_clone not implemented for '{name}', audio_url={audio_url}")
            else:
                summary = '🎙️ 请提供音频URL来克隆声音，或到「🎙️ 多媒体」Tab上传音频文件。'

        elif intent == 'video':
            title = args.get('title', _('Digital human video'))
            text = args.get('text', message)
            voice_id = args.get('voice_id', '')
            image_url = args.get('image_url', '')
            summary = f'🎬 视频创作："{title}"\n请到「🎙️ 多媒体 → 🎬 视频创作」Tab选择声音后提交。'

        elif intent == 'cms':
            title = args.get('title', '')
            summary = f'📝 CMS文章：{title or "待定"}\\n请切换到「📝 文章」Tab编辑并提交。'

        elif intent == 'clean':
            content = args.get('content', message)
            try:
                sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'auth-center', 'routes'))
                from cleaner_agent import process_clean_content
                res = process_clean_content(content, admin_id=admin.get('user_id', 0))
                if res['success']:
                    title = res.get('title', '')
                    cat = res.get('category', '')
                    kb_id = res.get('kb_id', '')
                    if kb_id == 'duplicate':
                        summary = _('⏭ Detected duplication, skipped:') + title
                    else:
                        summary = _('✅ Cleaning Complete:') + title + ' (' + cat + ')\n知识库ID: ' + kb_id
                    actions.append({'type': 'info', 'text': _('Written to knowledge_blocks, takes effect automatically for mini program/customer service')})
                else:
                    summary = _('❌ Cleaning Failed:') + res.get('error', _('Unknown error'))
            except Exception as e:
                import traceback
                traceback.print_exc()
                summary = _('❌ Cleaning Exception:') + str(e)

        elif intent == 'supply_chain':
            # 供应链/商城操作 → 交给 Master Agent 通过 Orchestrator 分配
            from agent_matrix.orchestrator import AgentOrchestrator
            orchestrator = AgentOrchestrator(models_module=_m())
            result = orchestrator.process_instruction(
                instruction=message,
                master_agent_id=master['id'],
                session_id=session_id,
                mode='fast',
                user_id=admin.get('user_id', 0)
            )
            # Write 层：对话结束自动提取知识（异步）
            orchestrator._on_task_complete(message, admin.get('user_id', 0), result)
            return _success({
                'session_id': session_id,
                'summary': result.get('summary', ''),
                'sub_task_results': result.get('sub_task_results', []),
                'actions': [],
                'status': result.get('status', 'ok'),
            })

        elif intent == 'ads':
            # 广告管理 → 直接调用插件 AI 工具
            import plugins.ads.ai_tools as ads_ai
            action = args.get('action', '')
            # 代码推断 action（LLM 不一定正确返回）
            if not action:
                msg_lower = message.lower()
                if any(w in msg_lower for w in ['创建','新增','生成','帮我创建','添加','create','add','new']):
                    action = 'create'
                elif any(w in msg_lower for w in ['更新','修改','编辑','update','edit','modify']):
                    action = 'update'
                elif any(w in msg_lower for w in ['删除','移除','delete','remove']):
                    action = 'delete'
                elif any(w in msg_lower for w in ['统计','数据','stats','statistics']):
                    action = 'stats'
                elif any(w in msg_lower for w in ['分析','analyze','insight']):
                    action = 'analyze'
                elif any(w in msg_lower for w in ['代码','snippet','渲染','render','模板']):
                    action = 'snippet'
                else:
                    action = 'list'
            result = {'success': False, 'error': _('Unknown ads operation')}
            try:
                if action == 'list':
                    result = ads_ai.list_ads(
                        site_key=args.get('site_key'),
                        position=args.get('position'),
                        active_only=args.get('active_only', False)
                    )
                elif action == 'create':
                    result = ads_ai.create_ad(args)
                elif action == 'update':
                    result = ads_ai.update_ad(
                        args.get('ad_id'),
                        args.get('updates', {})
                    )
                elif action == 'delete':
                    result = ads_ai.delete_ad(args.get('ad_id'))
                elif action == 'stats':
                    result = ads_ai.get_stats(
                        ad_id=args.get('ad_id'),
                        site_key=args.get('site_key'),
                        days=int(args.get('days', 7))
                    )
                elif action == 'analyze':
                    result = ads_ai.analyze_ads(days=int(args.get('days', 7)))
                elif action == 'snippet':
                    result = ads_ai.generate_render_snippet(
                        position=args.get('position', 'sidebar'),
                        page=args.get('page', '*'),
                        site_key=args.get('site_key', 'default'),
                        zone_id=args.get('zone_id')
                    )
            except Exception as ads_err:
                result = {'success': False, 'error': str(ads_err)}

            if result['success']:
                data = result.get('data')
                if isinstance(data, str):
                    summary = data
                elif isinstance(data, dict) and 'id' in data:
                    summary = f'✅ Operation successful, ID: {data["id"]}'
                elif isinstance(data, list):
                    summary = '\n'.join(str(x) for x in data[:30])
                else:
                    summary = _('✅ Operation successful')
            else:
                summary = f'❌ Operation Failed: {result.get("error", "Unknown Error")}'

            return _success({
                'session_id': session_id,
                'summary': summary,
                'sub_task_results': [],
                'actions': [],
                'status': 'ok',
                'intent': intent,
            })

        elif intent == 'site_build':
            # AI 智能建站 → 通过 Site Builder 引擎处理
            action = args.get('action', 'preview')
            prompt_identifier = args.get('prompt_identifier', '')
            try:
                from plugins.site_builder.models import get_prompt as _get_prompt, list_prompts as _list_prompts
                from plugins.site_builder.engine import SiteBuilderEngine
                engine = SiteBuilderEngine()

                # 获取提示词模板
                if prompt_identifier:
                    prompt_template = _get_prompt(prompt_identifier)
                else:
                    prompts = _list_prompts(active_only=True)
                    prompt_template = prompts[0] if prompts else None

                if not prompt_template:
                    summary = '❌ 没有可用的行业提示词模板，请先在「AI 建站」中创建模板。'
                    return _success({
                        'session_id': session_id,
                        'summary': summary,
                        'sub_task_results': [],
                        'actions': [],
                        'status': 'ok',
                        'intent': intent,
                    })

                if action == 'execute':
                    # 执行建站：需要前端传来的 plan 数据
                    summary = '✅ 请切换到「AI 智能建站」页面，在方案预览中点击「确认执行」按钮来启动建站流程。'
                    actions.append({
                        'type': 'navigate',
                        'text': _('Go to AI Smart Website Builder'),
                        'url': '/admin/site-builder'
                    })
                elif action == 'modify':
                    # 最小化修改
                    modify_result = engine.modify_block(message)
                    if modify_result.get('success'):
                        summary = f'✅ Modified: {modify_result.get("old_value", "")} → {modify_result.get("new_value", "")}'
                    else:
                        summary = f'❌ {modify_result.get("error", "Cannot Locate Block to Modify")}'
                else:
                    # 默认：生成方案预览
                    parsed = engine.parse_requirement(prompt_template, message)
                    plan = engine.generate_plan(prompt_template, parsed, message)
                    summary = plan.get('summary', _('Plan Generated'))

                    # 返回方案数据供前端展示
                    actions.append({
                        'type': 'site_build_plan',
                        'plan': plan,
                        'parsed': parsed,
                        'prompt_id': prompt_template.get('id', 0),
                    })

            except Exception as e:
                import traceback
                traceback.print_exc()
                summary = f'❌ Website Construction Failed: {e}'

        else:
            # 普通对话
            from agent_matrix.orchestrator import AgentOrchestrator
            orchestrator = AgentOrchestrator(models_module=_m())
            result = orchestrator.process_instruction(
                instruction=message,
                master_agent_id=master['id'],
                session_id=session_id,
                mode='fast',
                user_id=admin.get('user_id', 0)
            )
            # Write 层：对话结束自动提取知识（异步）
            orchestrator._on_task_complete(message, admin.get('user_id', 0), result)
            return _success({
                'session_id': session_id,
                'summary': result.get('summary', ''),
                'sub_task_results': result.get('sub_task_results', []),
                'actions': [],
                'status': result.get('status', 'ok'),
            })

    except Exception as e:
        import traceback
        traceback.print_exc()
        summary = f'❌ Execution Error: {e}'

    return _success({
        'session_id': session_id,
        'summary': summary,
        'sub_task_results': [],
        'actions': actions,
        'status': 'ok',
        'intent': intent,
    })


def _generate_ppt_file(topic, pages=10, style=_('Dark Tech Style, 16:9')):
    """生成PPT文件，返回下载文件名"""
    import os, json, time, uuid
    try:
        from agent_matrix.engine import UnifiedLLM
        agents = _m().list_agents(role_type='master', active_only=True)
        if not agents:
            return None
        engine = UnifiedLLM(agents[0])

        # AI 生成大纲
        outline_prompt = f"""你是一个PPT大纲生成器。主题：{topic}，需要{pages}页，风格：{style}。
返回JSON格式：
{{"title":"...","subtitle":"...","slides":[{{"title":"...","content":[_("Key Point 1"),_("Key Point 2"),...],"layout":"bullet"}}]}}
每页3-5个要点，第1页是封面，最后1页是总结。只返回JSON。"""
        raw = engine.chat([{"role": "user", "content": outline_prompt}], temperature=0.7, max_tokens=4000)
        import re
        m = re.search(r'\{[\s\S]*\}', raw)
        outline = json.loads(m.group(0)) if m else {'slides': []}

        # 生成PPTX
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        BG = RGBColor(0x0D, 0x11, 0x17)
        ACCENT = RGBColor(0x00, 0xD4, 0xAA)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        MUTED = RGBColor(0x88, 0x99, 0xAA)

        for i, s in enumerate(outline.get('slides', [])):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background
            bg.fill.solid()
            bg.fill.fore_color.rgb = BG

            s_title = s.get('title', '')
            s_content = s.get('content', [])

            # Accent line
            line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.2) if i > 0 else Inches(2.8), Inches(1.2), Pt(3))
            line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

            # Title
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4) if i > 0 else Inches(2.2), Inches(11.7), Inches(0.8))
            tf = txBox.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.text = s_title
            p.font.size = Pt(36) if i == 0 else Pt(28)
            p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = 'Arial'

            if i == 0:
                sub = outline.get('subtitle', '')
                if sub:
                    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.6))
                    tf2 = txBox2.text_frame
                    p2 = tf2.paragraphs[0]; p2.text = sub
                    p2.font.size = Pt(18); p2.font.color.rgb = MUTED; p2.font.name = 'Arial'
            else:
                for j, item in enumerate(s_content):
                    y = 2.6 + j * 0.6
                    txBox3 = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.0), Inches(0.5))
                    tf3 = txBox3.text_frame; tf3.word_wrap = True
                    p3 = tf3.paragraphs[0]; p3.text = f"▸ {item}"
                    p3.font.size = Pt(16); p3.font.color.rgb = WHITE; p3.font.name = 'Arial'

        # 保存
        media_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'data', 'media', 'temp')
        os.makedirs(media_dir, exist_ok=True)
        filename = f"ppt_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(media_dir, filename)
        prs.save(filepath)
        return filename
    except Exception as e:
        import traceback
        traceback.print_exc()
        return None


@agent_matrix_bp.route('/md-preview/<filename>')
def agent_md_preview(filename):
    """渲染 Markdown 文件为 HTML 预览"""
    admin, err = _require_admin()
    if err: return err

    import markdown as _md
    import bleach as _bleach
    media_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'data', 'media', 'temp')
    filename = os.path.basename(filename)
    fp = os.path.join(media_dir, filename)
    if not os.path.exists(fp):
        return jsonify({'success': False, 'error': _('File does not exist')}), 404
    try:
        with open(fp, 'r', encoding='utf-8') as f:
            md_content = f.read()
        html = _md.markdown(md_content, extensions=['fenced_code', 'tables', 'nl2br'])
        # 白名单清洗 HTML，防止存储型 XSS
        allowed_tags = ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'ul', 'ol', 'li',
                        'a', 'strong', 'em', 'code', 'pre', 'blockquote', 'table',
                        'thead', 'tbody', 'tr', 'th', 'td', 'br', 'hr', 'img']
        allowed_attrs = {'a': ['href', 'title'], 'img': ['src', 'alt', 'title']}
        html = _bleach.clean(html, tags=allowed_tags, attributes=allowed_attrs)
        return jsonify({'success': True, 'html': html, 'filename': filename})
    except Exception as e:
        import logging
        logging.getLogger(__name__).error(f"md-preview failed for '{filename}': {e}")
        return jsonify({'success': False, 'error': _('File preview failed')}), 500


@agent_matrix_bp.route('/chat/history', methods=['GET'])
def chat_history():
    admin, err = _require_admin()
    if err: return err

    sessions = _m().list_sessions(limit=30)
    return _success(sessions)


@agent_matrix_bp.route('/chat/<session_id>', methods=['GET'])
def chat_detail(session_id):
    admin, err = _require_admin()
    if err: return err

    messages = _m().get_conversation(session_id)
    return _success(messages)


@agent_matrix_bp.route('/chat/<session_id>/clear', methods=['POST'])
def clear_chat(session_id):
    admin, err = _require_admin()
    if err: return err

    with _m().get_db() as conn:
        conn.execute("DELETE FROM agent_conversations WHERE session_id=%s", (session_id,))
        conn.commit()
    return _success(None, _('Session has been cleared'))


@agent_matrix_bp.route('/chat/search', methods=['GET'])
def chat_search():
    """全文检索会话"""
    admin, err = _require_admin()
    if err: return err

    keyword = request.args.get('q', '').strip()
    if not keyword or len(keyword) < 2:
        return _error(_('Keyword must be at least 2 characters'))

    results = _m().search_conversations(keyword)
    return _success(results)


@agent_matrix_bp.route('/chat/batch-delete', methods=['POST'])
def chat_batch_delete():
    """批量删除会话"""
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    session_ids = data.get('session_ids', [])
    if not session_ids or not isinstance(session_ids, list):
        return _error('请提供 session_ids 列表')

    count = _m().batch_delete_sessions(session_ids)
    return _success({'deleted': count}, f'Deleted {count} records')


# ============================================================
# 4. 调度（直接分发任务给指定 Agent）
# ============================================================

@agent_matrix_bp.route('/dispatch', methods=['POST'])
def dispatch_task():
    """
    直接下发任务给指定 Sub Agent（绕过 Master Agent 分解）

    请求: {
        "target_agent_id": 2,
        "title": _("Create CMS Article"),
        "description": "...",
        "input_data": {...}
    }

    媒体 Agent（domain='media'）特殊处理：
    { "target_agent_id": 10, "action": "voice_clone", "params": {...} }
    """
    admin, err = _require_admin()
    if err: return err

    ai_err = _check_ai_access()
    if ai_err: return ai_err

    data = request.get_json(force=True) or {}
    target_id = data.get('target_agent_id')

    agent_config = _m().get_agent(target_id)
    if not agent_config:
        return _error(_('Target Agent does not exist'), 404)
    if not agent_config['is_active']:
        return _error(_('Target Agent is disabled'))

    # ── Image Agent special path: direct API call ──
    if agent_config.get('domain') == 'image':
        action = data.get('action', 'generate_image')
        params = data.get('params', {})

        prompt_text = params.get('prompt', params.get('text', ''))
        title = params.get('title', '')
        from services.ai_content_generator import generate_image as gen_img
        from services.ai_content_generator import generate_cover_image
        if title:
            img_url = generate_cover_image(title, prompt_text or title)
        else:
            img_url = gen_img(prompt_text or _('Illustration'))
        result = {'image_url': img_url, 'status': 'completed' if img_url else 'failed'}

        # 记录 token（按次计费）
        import threading
        from agent_matrix.engine import _log_token_usage
        media_user = data.get('user_id') if isinstance(data.get('user_id'), int) else None
        threading.Thread(target=_log_token_usage, args=(
            target_id, agent_config.get('name', 'image'),
            agent_config.get('model_name', ''), agent_config.get('provider', ''),
            0, 0, 1, action, 'image', media_user
        ), daemon=True).start()

        return _success(result)

    # ── 标准文本 Agent 路径 ──
    if not target_id:
        return _error(_('Target_agent_id must'))

    # 加载 prompt（动态解析）
    from agent_matrix.orchestrator import AgentOrchestrator
    orch = AgentOrchestrator(models_module=_m())
    task_ctx = {
        'domain': agent_config.get('domain', 'general'),
        'task_type': data.get('task_type', data.get('action', 'execute')),
        'mode': data.get('mode', ''),
        'user_query': data.get('description', '') or data.get('title', ''),
    }
    prompt = orch._resolve_prompt(agent_config, task_ctx)
    if prompt:
        agent_config['system_prompt'] = prompt

    from agent_matrix.agent_runner import AgentRunner
    runner = AgentRunner(agent_config, db_models=_m())

    task_id = _m().create_task({
        'source_agent_id': 0,
        'target_agent_id': target_id,
        'task_type': 'execute',
        'title': data.get('title', ''),
        'description': data.get('description', ''),
        'input_data': data.get('input_data', {}),
        'expected_output': data.get('expected_output', {}),
        'target_module': agent_config.get('domain', ''),
        'max_retries': 2,
    })
    _m().update_task_status(task_id, 'running')

    result = runner.execute({
        'task_id': task_id,
        'title': data.get('title', ''),
        'description': data.get('description', ''),
        'input_data': data.get('input_data', {}),
        'expected_output': data.get('expected_output', {}),
        'max_retries': 2,
    })

    if result['status'] == 'completed':
        _m().update_task_status(task_id, 'completed',
                                result_data=result.get('response', ''),
                                confidence=result.get('confidence', 0.9))
    else:
        _m().update_task_status(task_id, 'failed',
                                error_message=result.get('response', ''))

    _m().update_agent_stats(target_id, success=(result['status'] == 'completed'))

    stats = _m().get_matrix_stats()
    return _success({
        'task_id': task_id,
        'result': result,
        'stats': stats
    })


# ============================================================
# 5. 提示词模板管理
# ============================================================

@agent_matrix_bp.route('/providers', methods=['GET'])
def list_providers():
    """返回可用提供商列表（从数据库 providers 表读取）"""
    admin, err = _require_admin()
    if err:
        return err
    with get_db() as conn:
        rows = conn.execute(
            "SELECT id, slug, name, description, is_active FROM providers "
            "WHERE is_active=1 ORDER BY id"
        ).fetchall()
        return _success([dict(r) for r in rows])


@agent_matrix_bp.route('/prompts', methods=['GET'])
def list_prompts():
    """数据库提示词列表（支持 ?type=&domain=&keyword= 筛选）"""
    admin, err = _require_admin()
    if err: return err

    prompt_type = request.args.get('type')
    domain = request.args.get('domain')
    keyword = request.args.get('keyword')
    try:
        rows = _m().list_prompts(
            prompt_type=prompt_type, domain=domain, keyword=keyword
        )
        return _success(rows)
    except Exception as e:
        return _error(f'List prompts failed: {e}', 500)


@agent_matrix_bp.route('/prompts/files', methods=['GET'])
def list_prompt_files():
    """列出可用 Prompt 文件模板（chat 编辑器下拉框使用）"""
    admin, err = _require_admin()
    if err: return err

    templates = [
        {"id": "custom", "name": "自定义提示词", "description": "手动编写完整的 System Prompt", "is_builtin": False},
    ]

    # 文件模板（保持 chat 编辑器兼容）
    prompt_dir = os.path.join(os.path.dirname(__file__), 'prompts')
    if os.path.exists(prompt_dir):
        for f in sorted(os.listdir(prompt_dir)):
            if f.endswith('.md'):
                filepath = os.path.join(prompt_dir, f)
                with open(filepath, 'r', encoding='utf-8') as fh:
                    first_line = fh.readline().strip().lstrip('#').strip()
                rel_path = f'prompts/{f}'
                name = f.replace('sub_', _('Sub Agent - ')).replace('_prompt.md', '').replace('master_', 'Master - ')
                if name == f.replace('.md', ''):
                    name = f.replace('.md', '')
                templates.append({
                    "id": rel_path,
                    "name": name,
                    "description": first_line or '',
                    "is_builtin": True,
                    "filename": f,
                })

    return _success(templates)


@agent_matrix_bp.route('/prompts/load', methods=['GET'])
def load_prompt_content():
    """加载指定 prompt 文件内容"""
    admin, err = _require_admin()
    if err: return err

    path = request.args.get('path', '')
    if not path:
        return _error(_('Path parameter must'))

    # Security: 仅允许 prompts/ 下文件，且用 realpath 防路径遍历（对齐 resolver 实现）
    prompt_dir = os.path.join(os.path.dirname(__file__), 'prompts')
    if not path.startswith('prompts/') or '..' in path:
        return _error(_('Invalid path'), 400)

    full_path = os.path.realpath(os.path.join(os.path.dirname(__file__), path))
    real_base = os.path.realpath(prompt_dir)
    if not full_path.startswith(real_base + os.sep):
        logging.getLogger(__name__).warning(f'Prompt 路径遍历尝试被拦截: {path}')
        return _error(_('Invalid path'), 400)

    if not os.path.exists(full_path):
        return _error(_('Prompt file does not exist'), 404)

    with open(full_path, 'r', encoding='utf-8') as f:
        content = f.read()

    return _success(content)


# ============================================================
# Prompt DB CRUD (Dynamic Prompt System) — 文档 12 节
# ============================================================

@agent_matrix_bp.route('/prompts', methods=['POST'])
def create_prompt_route():
    """创建新提示词"""
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    if not data.get('name'):
        return _error(_('Name is required'))
    if not data.get('content'):
        return _error(_('Content is required'))
    # prompt_type 枚举校验（对齐 DB CHECK 约束，提前返回 400）
    prompt_type = data.get('prompt_type', 'system')
    if prompt_type not in ('system', 'scene', 'tool', 'rule', 'composite'):
        return _error(_('Invalid prompt_type'), 400)
    try:
        new_id = _m().create_prompt(data)
        return _success({'id': new_id})
    except Exception as e:
        return _error(f'Create prompt failed: {e}', 500)


@agent_matrix_bp.route('/prompts/<int:pid>', methods=['GET'])
def get_prompt_route(pid):
    """提示词详情"""
    admin, err = _require_admin()
    if err: return err

    row = _m().get_prompt(pid)
    if not row:
        return _error(_('Prompt does not exist'), 404)
    return _success(row)


@agent_matrix_bp.route('/prompts/<int:pid>', methods=['PUT'])
def update_prompt_route(pid):
    """更新提示词（content 变更自动 version+1）"""
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    try:
        result = _m().update_prompt(pid, data)
        if result is None:
            return _error(_('Prompt does not exist'), 404)
        return _success({'id': result})
    except Exception as e:
        return _error(f'Update prompt failed: {e}', 500)


@agent_matrix_bp.route('/prompts/<int:pid>', methods=['DELETE'])
def delete_prompt_route(pid):
    """软删除提示词（is_active=false）"""
    admin, err = _require_admin()
    if err: return err

    row = _m().get_prompt(pid)
    if not row:
        return _error(_('Prompt does not exist'), 404)
    _m().delete_prompt(pid)
    return _success()


@agent_matrix_bp.route('/prompts/<int:pid>/versions', methods=['GET'])
def prompt_versions_route(pid):
    """同 slug 的版本历史"""
    admin, err = _require_admin()
    if err: return err

    row = _m().get_prompt(pid)
    if not row:
        return _error(_('Prompt does not exist'), 404)
    versions = _m().get_prompt_versions(row['slug'])
    return _success(versions)


@agent_matrix_bp.route('/prompts/<int:pid>/test', methods=['POST'])
def test_prompt_route(pid):
    """用指定 Prompt 回答测试问题（验证组装效果）"""
    admin, err = _require_admin()
    if err: return err

    ai_err = _check_ai_access()
    if ai_err: return ai_err

    row = _m().get_prompt(pid)
    if not row:
        return _error(_('Prompt does not exist'), 404)

    data = request.get_json(force=True) or {}
    question = data.get('question', '')
    if not question:
        return _error(_('Question is required'))

    # 用测试 Agent 的模型配置（默认第一个 active sub agent）
    agent_config = None
    agent_id = data.get('agent_id')
    if agent_id:
        agent_config = _m().get_agent(agent_id)
    if not agent_config:
        agents = _m().list_agents(role_type='sub', active_only=True)
        agent_config = agents[0] if agents else None
    if not agent_config:
        return _error(_('No available Agent'), 500)

    agent_config['system_prompt'] = row['content']

    from agent_matrix.engine import UnifiedLLM
    engine = UnifiedLLM(agent_config)
    if not engine.is_ready():
        return _error(_('AI Engine Not Ready (Check API Key)'), 500)

    try:
        response = engine.ask(question, temperature=0.3)
        return _success({'prompt_id': pid, 'response': response})
    except Exception as e:
        return _error(f'Test failed: {e}', 500)


@agent_matrix_bp.route('/agents/<int:aid>/bindings', methods=['GET'])
def list_agent_bindings_route(aid):
    """列出指定 Agent 的 Prompt 绑定关系"""
    admin, err = _require_admin()
    if err: return err

    if not _m().get_agent(aid):
        return _error(_('Agent does not exist'), 404)
    try:
        rows = _m().list_bindings(agent_id=aid)
        return _success(rows)
    except Exception as e:
        return _error(f'List bindings failed: {e}', 500)


@agent_matrix_bp.route('/agents/<int:aid>/bind-prompt', methods=['POST'])
def bind_prompt_route(aid):
    """为 Agent 绑定 Prompt"""
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    prompt_id = data.get('prompt_id')
    if not prompt_id:
        return _error(_('prompt_id is required'))

    if not _m().get_agent(aid):
        return _error(_('Agent does not exist'), 404)
    if not _m().get_prompt(prompt_id):
        return _error(_('Prompt does not exist'), 404)

    binding_type = data.get('binding_type', 'default')
    if binding_type not in ('default', 'scene', 'override', 'mode'):
        return _error(_('Invalid binding_type'), 400)

    _m().create_binding(
        aid, prompt_id,
        binding_type=binding_type,
        condition=data.get('condition', ''),
        priority=data.get('priority', 0),
    )
    return _success()


@agent_matrix_bp.route('/agents/<int:aid>/bind-prompt/<int:bid>', methods=['DELETE'])
def unbind_prompt_route(aid, bid):
    """解绑 Agent-Prompt（校验绑定归属该 Agent）"""
    admin, err = _require_admin()
    if err: return err

    if not _m().get_agent(aid):
        return _error(_('Agent does not exist'), 404)
    deleted = _m().delete_binding(bid, agent_id=aid)
    if not deleted:
        return _error(_('Binding does not exist for this agent'), 404)
    return _success()


# ============================================================
# 5B. AI 服务能力管理
# ============================================================

@agent_matrix_bp.route('/ai-services', methods=['GET'])
def list_ai_services():
    """列出所有 AI 能力及其在各 Agent 中的使用情况"""
    admin, err = _require_admin()
    if err: return err

    services = []
    agents = _m().list_agents(active_only=True)

    # Read system config keys
    keys = {}
    with _m().get_db() as conn:
        rows = conn.execute(
            "SELECT key, value FROM system_config WHERE key IN ('dashscope_text_key','dashscope_api_key')"
        ).fetchall()
        keys = {r['key']: r['value'] for r in rows}

    # Key status helper
    def key_status(ref_name):
        val = keys.get(ref_name, '')
        return 'configured' if val else 'missing'

    # 1. Qwen 文本生成
    qwen_agents = [a for a in agents if a.get('api_key_ref') == 'dashscope_text_key'
                   or a.get('provider') == 'dashscope']
    services.append({
        "id": "qwen_text",
        "name": _("Text generation (Qwen-turbo)"),
        "type": "text_generation",
        "provider": "DashScope",
        "key_ref": "dashscope_text_key",
        "key_status": key_status('dashscope_text_key'),
        "models": ["qwen-turbo", "qwen-max", "qwen-plus"],
        "default_model": "qwen-turbo",
        "used_by_agents": [{"id": a['id'], "name": a['name'], "domain": a['domain']} for a in qwen_agents],
        "endpoints": [
            _("/admin/content-factory/ai-format — AI Formatting"),
            _("/admin/content-factory/process — AI Content Processing"),
        ],
    })

    # 2. 通义万相文生图
    image_agents = [a for a in agents if a.get('domain') in ('cms', 'content-factory', 'image')
                    or 'image_gen' in str(a.get('capabilities', ''))]
    services.append({
        "id": "wanx_image",
        "name": _("Text to Image (Tongyi Wan2.7)"),
        "type": "image_generation",
        "provider": "DashScope",
        "key_ref": "dashscope_api_key",
        "key_status": key_status('dashscope_api_key'),
        "models": ["wanx2.1-t2i-turbo"],
        "default_model": "wanx2.1-t2i-turbo",
        "used_by_agents": [{"id": a['id'], "name": a['name'], "domain": a['domain']} for a in image_agents],
        "endpoints": [
            _("/admin/content-factory/ai-cover — AI Image"),
            _("/admin/social/generate-image — Social Media Image"),
        ],
    })

    # 3. Agent Matrix 对话引擎 (Master Agent)
    master = [a for a in agents if a['role_type'] == 'master']
    services.append({
        "id": "matrix_chat",
        "name": _("Matrix Dialogue Engine (Master Agent)"),
        "type": "orchestration",
        "provider": master[0]['provider'] if master else '—',
        "key_ref": master[0].get('api_key_ref', '') if master else '',
        "key_status": key_status(master[0].get('api_key_ref', '')) if master else 'missing',
        "models": [master[0]['model_name']] if master else [],
        "used_by_agents": [{"id": a['id'], "name": a['name'], "domain": a['domain']} for a in agents if a['role_type'] == 'sub'],
        "endpoints": [
            _("POST /admin/agent-matrix/chat — Send Instructions to Master Agent"),
        ],
    })

    # 4. Trademind 聊天窗口 (独立)
    services.append({
        "id": "trademind_chat",
        "name": _("💬 Chat Window (TradeMind Support)"),
        "type": "chat",
        "provider": "DashScope",
        "key_ref": "dashscope_text_key",
        "key_status": key_status('dashscope_text_key'),
        "models": ["qwen-turbo"],
        "used_by_agents": [],
        "endpoints": [
            _("Located in trademind/chatbot.py - Standalone service, not yet integrated into the matrix"),
        ],
        "note": _("This service runs independently as TradeMind (8081), it is recommended to migrate to Kai Assistant for unified management later")
    })

    return _success(services)


# ============================================================
# 6. 聊天窗口迁移 — SSE 流式 + 知识库管理
# ============================================================

@agent_matrix_bp.route('/chat/stream', methods=['POST'])
def chat_stream_sse():
    """
    SSE 流式聊天接口（管理员测试用）
    请求: { message, history, agent_id (可选, 默认用 Kai Assistant) }
    返回: text/event-stream
    """
    admin, err = _require_admin()
    if err: return err

    ai_err = _check_ai_access()
    if ai_err: return ai_err

    data = request.get_json(force=True) or {}
    message = (data.get('message') or '').strip()
    history = data.get('history', [])
    agent_id = data.get('agent_id')

    if not message:
        return jsonify({'error': _('Message is required')}), 400

    # 获取 Agent 配置
    if agent_id:
        agent_config = _m().get_agent(agent_id)
    else:
        # 默认用 Advisor Agent（chatbot 域）或第一个 sub agent
        agents = _m().list_agents(role_type='sub', active_only=True)
        chatbot = [a for a in agents if a.get('domain') == 'chatbot']
        agent_config = chatbot[0] if chatbot else (agents[0] if agents else None)

    if not agent_config:
        return jsonify({'error': _('No available Agent')}), 500

    # 加载 system_prompt（动态解析）
    from agent_matrix.orchestrator import AgentOrchestrator
    orch = AgentOrchestrator(models_module=_m())
    base_prompt = orch._resolve_prompt(agent_config, {
        'domain': agent_config.get('domain', 'general'),
        'task_type': 'chat',
        'mode': data.get('mode', ''),
        'user_query': message[:200],
    })

    # 加载知识库（实时 RAG 检索 + 静态兜底）
    knowledge = _inject_knowledge(user_message=message)
    # 兜底：保留 system_config.chatbot_knowledge_base 作为静态基础知识
    static_knowledge = ''
    with _m().get_db() as conn:
        row = conn.execute(
            "SELECT value FROM system_config WHERE key='chatbot_knowledge_base'"
        ).fetchone()
        if row and row['value']:
            static_knowledge = row['value']

    # 构建完整的 system prompt（先静态知识，后动态知识）
    full_system = base_prompt
    if static_knowledge:
        full_system += f"\n\n=== 基础知识库 ===\n{static_knowledge}\n=== 基础知识库结束 ==="
    if knowledge:
        full_system += knowledge

    agent_config['system_prompt'] = full_system

    # 创建 AI 引擎
    from agent_matrix.engine import UnifiedLLM
    engine = UnifiedLLM(agent_config)

    if not engine.is_ready():
        def err_gen():
            yield f"data: {json.dumps(_('AI Engine Not Ready (Check API Key)'))}\n\n"
            yield "data: [DONE]\n\n"
        from flask import Response, stream_with_context
        return Response(
            stream_with_context(err_gen()),
            mimetype='text/event-stream',
            headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'}
        )

    # 构建 messages
    messages = [{"role": "system", "content": full_system}]
    for h in (history or []):
        if h.get('role') in ('user', 'assistant'):
            messages.append({"role": h['role'], "content": h['content']})
    messages.append({"role": "user", "content": message})

    from flask import Response, stream_with_context

    def generate():
        # 先发 role 标识
        yield "data: {\"role\":\"assistant\"}\n\n"
        for chunk in engine.chat_stream(messages):
            yield f"data: {json.dumps(chunk)}\n\n"
        yield "data: [DONE]\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no'}
    )


# DEPRECATED: 此端点操作的是 system_config.chatbot_knowledge_base（旧版文本知识库），
# 非 knowledge_blocks 表（RAG 知识库）。RAG 知识库管理请使用 /api/v1/knowledge/* 端点。
# 保留此端点以兼容 ai_chat.html 中的 matEditKnowledge() 功能。
@agent_matrix_bp.route('/chat/knowledge', methods=['GET'])
def get_knowledge_base():
    """获取聊天知识库内容（旧版）"""
    admin, err = _require_admin()
    if err: return err

    content = ''
    with _m().get_db() as conn:
        row = conn.execute(
            "SELECT value FROM system_config WHERE key='chatbot_knowledge_base'"
        ).fetchone()
        if row:
            content = row['value']
    return _success(content)


@agent_matrix_bp.route('/chat/knowledge', methods=['PUT'])
def update_knowledge_base():
    """更新聊天知识库"""
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    content = data.get('content', '')

    with _m().get_db() as conn:
        conn.execute("""
            INSERT INTO system_config (key, value, description, updated_at)
            VALUES ('chatbot_knowledge_base', %s, _('Chatbot Knowledge Base'), NOW())
            ON CONFLICT(key) DO UPDATE SET value=excluded.value, updated_at=NOW()
        """, (content,))
        conn.commit()

    return _success(None, _('Knowledge base has been updated'))


# ============================================================
# 7. 统计与监控
# ============================================================

@agent_matrix_bp.route('/stats', methods=['GET'])
def matrix_stats():
    admin, err = _require_admin()
    if err: return err

    stats = _m().get_matrix_stats()
    return _success(stats)


@agent_matrix_bp.route('/dashboard', methods=['GET'])
def matrix_dashboard():
    admin, err = _require_admin()
    if err: return err

    stats = _m().get_matrix_stats()
    recent = _m().get_recent_tasks(limit=10)
    agents = _m().list_agents(active_only=True)

    return _success({
        'stats': stats,
        'recent_tasks': recent,
        'active_agents': agents,
    })


@agent_matrix_bp.route('/health', methods=['GET'])
def health_check():
    """所有 Agent 健康检查"""
    admin, err = _require_admin()
    if err: return err

    agents = _m().list_agents(active_only=True)
    results = []
    for a in agents:
        from agent_matrix.engine import UnifiedLLM
        engine = UnifiedLLM(a)
        results.append({
            'id': a['id'],
            'name': a['name'],
            'role_type': a['role_type'],
            'provider': a['provider'],
            'model': a['model_name'],
            'ready': engine.is_ready(),
        })

    return _success({
        'total': len(results),
        'ready': sum(1 for r in results if r['ready']),
        'agents': results
    })


# ============================================================
# 4b. Token 用量统计 (2026-05-16)
# ============================================================

@agent_matrix_bp.route('/token-stats', methods=['GET'])
def token_stats():
    """Token 消耗统计：今日/本周/本月/累计 + Agent×模型交叉 + 维度汇总 + 费用估算"""
    admin, err = _require_admin()
    if err: return err

    period   = request.args.get('period', 'today')
    dim      = request.args.get('dimension', '').strip()  # text/voice/video/image, 空=全部
    # 白名单校验，防止 SQL 注入
    ALLOWED_DIMENSIONS = {'text', 'voice', 'video', 'image'}

    date_where_base = {
        'today':  "created_at::date = CURRENT_DATE",
        'week':   "created_at::date >= CURRENT_DATE - INTERVAL '6 days'",
        'month':  "created_at::date >= DATE_TRUNC('month', CURRENT_DATE)::DATE",
        'all':    "1=1"
    }.get(period, "created_at::date = CURRENT_DATE")
    date_where = date_where_base
    params = []
    if dim and dim in ALLOWED_DIMENSIONS:
        date_where += " AND dimension = %s"
        params.append(dim)

    date_where_t = {
        'today':  "t.created_at::date = CURRENT_DATE",
        'week':   "t.created_at::date >= CURRENT_DATE - INTERVAL '6 days'",
        'month':  "t.created_at::date >= DATE_TRUNC('month', CURRENT_DATE)::DATE",
        'all':    "1=1"
    }.get(period, "t.created_at::date = CURRENT_DATE")
    if dim and dim in ALLOWED_DIMENSIONS:
        date_where_t += " AND t.dimension = %s"

    # 费用单价（后续可从 system_config 读取）
    pricing = {
        'text_per_1k': 0.003,    # ¥ / 1000 tokens
        'image_per_call': 0.05,   # ¥ / 次
    }

    try:
        with get_db() as conn:
            # ── 全局汇总 ──
            total_row = conn.execute(f"""
                SELECT COALESCE(SUM(total_tokens),0) AS total,
                       COALESCE(SUM(prompt_tokens),0) AS prompt,
                       COALESCE(SUM(completion_tokens),0) AS completion,
                       COUNT(*) AS calls
                FROM agent_token_logs WHERE {date_where}
            """, params).fetchone()

            # ── 按维度汇总 ──
            by_dim_rows = conn.execute(f"""
                SELECT dimension,
                       COALESCE(SUM(total_tokens),0) AS total,
                       COUNT(*) AS calls
                FROM agent_token_logs WHERE {date_where_base}
                GROUP BY dimension
                ORDER BY total DESC
            """).fetchall()

            # ── Agent×模型交叉 ──
            am_rows = conn.execute(f"""
                SELECT t.agent_id, t.agent_name,
                       COALESCE(t.model_name, '') AS model_name,
                       COALESCE(t.provider, '') AS provider,
                       COALESCE(t.dimension, 'text') AS dimension,
                       COALESCE(SUM(t.total_tokens),0) AS total,
                       COALESCE(SUM(t.prompt_tokens),0) AS prompt,
                       COALESCE(SUM(t.completion_tokens),0) AS completion,
                       COUNT(*) AS calls
                FROM agent_token_logs t
                WHERE {date_where_t}
                GROUP BY t.agent_id, t.agent_name, t.model_name, t.provider, t.dimension
                ORDER BY total DESC
            """, params).fetchall()

            # ── 按用户汇总 ──
            unknown_label = _('Unknown user')
            user_rows = conn.execute(f"""
                SELECT t.user_id,
                       COALESCE(u.username, u.phone, '{unknown_label}') AS username,
                       t.agent_id, t.agent_name,
                       COALESCE(t.model_name, '') AS model_name,
                       COALESCE(t.dimension, 'text') AS dimension,
                       COALESCE(SUM(t.total_tokens),0) AS total,
                       COUNT(*) AS calls
                FROM agent_token_logs t
                LEFT JOIN users u ON u.id = t.user_id
                WHERE {date_where_t}
                GROUP BY t.user_id, t.agent_id, t.agent_name, t.model_name, t.dimension, COALESCE(u.username, u.phone, '{unknown_label}')
                ORDER BY total DESC
                LIMIT 50
            """, params).fetchall()

            # ── 费用估算 ──
            text_tokens  = float(sum(int(r['total']) if r.get('dimension') == 'text' else 0 for r in by_dim_rows))
            image_calls  = int(sum(r['calls'] for r in by_dim_rows if r['dimension'] == 'image'))
            cost_est = (
                text_tokens / 1000 * pricing['text_per_1k'] +
                image_calls * pricing['image_per_call']
            )

            # ── 今日预警数据 ──
            today_total = conn.execute("""
                SELECT COALESCE(SUM(total_tokens),0) AS total
                FROM agent_token_logs WHERE created_at::date = CURRENT_DATE
            """).fetchone()
            today_by_agent = conn.execute("""
                SELECT agent_id, agent_name, COALESCE(SUM(total_tokens),0) AS total
                FROM agent_token_logs
                WHERE created_at::date = CURRENT_DATE
                GROUP BY agent_id, agent_name
            """).fetchall()

        total_row_d = dict(total_row)
        for k in total_row_d:
            total_row_d[k] = float(total_row_d[k]) if hasattr(total_row_d[k], 'real') else total_row_d[k]
        by_dim_d = []
        for r in by_dim_rows:
            d = dict(r)
            d['total'] = float(d['total'])
            d['calls'] = int(d['calls'])
            by_dim_d.append(d)
        today_total_val = float(today_total['total']) if today_total else 0
        return _success({
            'period': period,
            'dimension': dim or 'all',
            'total': total_row_d,
            'by_dimension': by_dim_d,
            'agent_models': [dict(r) for r in am_rows],
            'users': [dict(r) for r in user_rows],
            'cost_estimate': round(cost_est, 2),
            'pricing': pricing,
            'today_matrix_total': today_total_val,
            'today_by_agent': [dict(r) for r in today_by_agent],
            'thresholds': {
                'agent_yellow': 200000,
                'agent_red': 500000,
                'matrix_red': 2000000
            }
        })
    except Exception as e:
        return _error(f'Token statistics query failed: {e}', 500)


@agent_matrix_bp.route('/token-logs', methods=['GET'])
def token_logs():
    """Token 消耗日志列表（分页，可选 agent_id 筛选）"""
    admin, err = _require_admin()
    if err: return err

    agent_id = request.args.get('agent_id', '').strip()
    limit = request.args.get('limit', 50, type=int)
    offset = request.args.get('offset', 0, type=int)

    try:
        with get_db() as conn:
            where = ''
            params = []
            if agent_id:
                where = 'WHERE agent_id = %s'
                params.append(int(agent_id))

            total = conn.execute(
                f'SELECT COUNT(*) as c FROM agent_token_logs {where}', params
            ).fetchone()['c']

            rows = conn.execute(f"""
                SELECT id, agent_id, agent_name, model_name, provider,
                       prompt_tokens, completion_tokens, total_tokens,
                       call_type, task_id, created_at
                FROM agent_token_logs {where}
                ORDER BY created_at DESC
                LIMIT %s OFFSET %s
            """, params + [limit, offset]).fetchall()

        return _success({
            'total': total,
            'limit': limit,
            'offset': offset,
            'logs': [dict(r) for r in rows]
        })
    except Exception as e:
        return _error(f'Token log query failed: {e}', 500)


# ============================================================
# 5. 图片生成 + 本地存储
# ============================================================

@agent_matrix_bp.route('/generate-image', methods=['POST'])
def generate_and_save_image():
    """生成图片并保存到本地"""
    admin, err = _require_admin()
    if err:
        return err

    data = request.get_json(force=True) or {}
    prompt = data.get('prompt', '').strip()
    title = data.get('title', '').strip()
    use_for_cover = data.get('cover', True)

    if not prompt and not title:
        return _error('请输入图片描述或文章标题')

    try:
        from services.ai_content_generator import generate_image as gen_img
        from services.ai_content_generator import generate_cover_image

        if use_for_cover and title:
            image_url = generate_cover_image(title, prompt or title)
        else:
            image_url = gen_img(prompt or f'Image: {title}')

        if not image_url:
            return _error(_('Picture generation failed: No image address returned'))

        # 下载图片到本地
        import uuid, urllib.request
        img_data = urllib.request.urlopen(image_url, timeout=30).read()

        ext = '.png'
        if 'jpg' in image_url or 'jpeg' in image_url:
            ext = '.jpg'
        elif 'webp' in image_url:
            ext = '.webp'

        filename = f'{uuid.uuid4().hex}{ext}'
        save_dir = TEMP_UPLOAD_DIR
        os.makedirs(save_dir, exist_ok=True)
        save_path = os.path.join(save_dir, filename)

        with open(save_path, 'wb') as f:
            f.write(img_data)

        local_url = f'/static/uploads/temp/{filename}'

        # 记录 token（图像生成按次计费）
        import threading
        from agent_matrix.engine import _log_token_usage
        threading.Thread(target=_log_token_usage, args=(
            None, 'Image Generator', 'dashscope-image', 'dashscope',
            0, 0, 1, 'generate_image', 'image'
        ), daemon=True).start()

        return _success({
            'image_url': local_url,
            'original_url': image_url,
            'filename': filename,
            'size': len(img_data),
        }, f'Picture saved to {local_url}')

    except Exception as e:
        import traceback
        traceback.print_exc()
        return _error(f'Failed to generate image: {str(e)}', 500)


# ============================================================
# 5b. 文件上传 / 下载
# ============================================================

TEMP_UPLOAD_DIR = os.path.join(BASE_DIR, '..', 'admin', 'static', 'uploads', 'temp')
TEMP_RETENTION_DAYS = 7

ALLOWED_EXTENSIONS = {
    'jpg', 'jpeg', 'png', 'gif', 'webp', 'bmp',
    'pdf', 'doc', 'docx', 'txt', 'md', 'csv', 'json',
    'zip',
}

def _allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


@agent_matrix_bp.route('/upload', methods=['POST'])
def upload_file():
    """上传文件到临时存储（7天有效）"""
    admin, err = _require_admin()
    if err:
        return err

    if 'file' not in request.files:
        return _error('请选择要上传的文件')

    file = request.files['file']
    if not file.filename:
        return _error(_('File name is empty'))

    if not _allowed_file(file.filename):
        return _error(f'Unsupported file type, allowed: {", ".join(sorted(ALLOWED_EXTENSIONS))}')

    try:
        import uuid, datetime

        orig_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else 'bin'
        filename = f'{uuid.uuid4().hex}.{orig_ext}'

        os.makedirs(TEMP_UPLOAD_DIR, exist_ok=True)
        save_path = os.path.join(TEMP_UPLOAD_DIR, filename)
        file.save(save_path)

        now = datetime.datetime.utcnow()
        expires = now + datetime.timedelta(days=TEMP_RETENTION_DAYS)
        size = os.path.getsize(save_path)

        return _success({
            'url': f'/static/uploads/temp/{filename}',
            'filename': filename,
            'original_name': file.filename,
            'size': size,
            'size_display': _fmt_size(size),
            'uploaded_at': now.isoformat(),
            'expires_at': expires.isoformat(),
        }, f'File uploaded, valid for {TEMP_RETENTION_DAYS} days')
    except Exception as e:
        import traceback
        traceback.print_exc()
        return _error(f'Upload failed: {str(e)}', 500)


@agent_matrix_bp.route('/download/<filename>', methods=['GET'])
def download_temp_file(filename):
    """下载临时文件（含过期检查）"""
    admin, err = _require_admin()
    if err: return err
    import datetime, time

    safe_name = os.path.basename(filename)
    filepath = os.path.join(TEMP_UPLOAD_DIR, safe_name)

    if not os.path.isfile(filepath):
        return _error(_('File does not exist'), 404)

    # 过期检查
    mtime = os.path.getmtime(filepath)
    age_days = (time.time() - mtime) / 86400
    if age_days > TEMP_RETENTION_DAYS:
        try:
            os.remove(filepath)
        except OSError:
            pass
        return _error(_('File has expired'), 410)

    return send_from_directory(TEMP_UPLOAD_DIR, safe_name, as_attachment=True)


def _fmt_size(size_bytes):
    """格式化文件大小"""
    if size_bytes < 1024:
        return f'{size_bytes}B'
    elif size_bytes < 1024 * 1024:
        return f'{size_bytes / 1024:.1f}KB'
    else:
        return f'{size_bytes / 1024 / 1024:.1f}MB'


# ============================================================
# 5c. TTS (Text-to-Speech) — Azure Neural TTS
# ============================================================

@agent_matrix_bp.route('/tts/synthesize', methods=['POST'])
def tts_synthesize():
    """Admin TTS endpoint: convert text to speech audio.

    Request JSON: {text: str, voice?: str}
    Response: audio/mpeg binary on success, JSON error on failure.
    Requires admin authentication.
    """
    admin, err = _require_admin()
    if err:
        return err

    data = request.get_json(force=True) or {}
    text = (data.get('text') or '').strip()
    if not text:
        return _error(_('Text is required'))
    if len(text) > 2000:
        return _error(_('Text too long (max 2000 characters)'))

    voice = data.get('voice', 'zh-CN-XiaoxiaoNeural')

    # 先尝试 Edge-TTS（免费，无需 Key）
    try:
        from services.tts_service import text_to_speech_bytes
        import asyncio
        audio_bytes = asyncio.run(text_to_speech_bytes(text, voice))
        if audio_bytes:
            from flask import Response
            return Response(
                audio_bytes,
                mimetype='audio/mpeg',
                headers={'Content-Length': str(len(audio_bytes))}
            )
    except Exception as e:
        pass

    # 回退 Azure TTS（需配置 Key）
    try:
        from agent_matrix.audio import AudioOutputProcessor
        processor = AudioOutputProcessor(provider='azure_tts', voice=voice)
        audio = processor.synthesize(text)
        if not audio:
            return _error(_('TTS synthesis failed'), 500)
        from flask import Response
        return Response(
            audio,
            mimetype='audio/mpeg',
            headers={'Content-Length': str(len(audio))}
        )
    except Exception as e:
        import logging
        logging.getLogger(__name__).error(
            '[TTS] Admin TTS failed: %s', e, exc_info=True
        )
        return _error(str(e), 500)


@agent_matrix_bp.route('/tts/voices', methods=['GET'])
def tts_list_voices():
    """列出可用语音 — 优先 Azure，不可用时回退 Edge-TTS"""
    admin, err = _require_admin()
    if err:
        return err

    locale = request.args.get('locale', '').strip()

    # 优先 Edge-TTS（免费，无需 Key）
    try:
        from services.tts_service import list_available_voices, VOICE_PRESETS
        import asyncio
        voices = asyncio.run(list_available_voices(locale or 'zh-CN'))
        return _success({
            'locale': locale or 'zh-CN',
            'count': len(voices),
            'voices': voices,
            'presets': VOICE_PRESETS,
        })
    except Exception:
        pass  # Edge 不可用，回退 Azure

    # Azure TTS 后备
    try:
        from agent_matrix.audio import AudioOutputProcessor
        processor = AudioOutputProcessor(provider='azure_tts')
        client = processor._get_client()
        if client:
            voices = client.list_voices(locale=locale)
            return _success({
                'locale': locale or 'all',
                'count': len(voices),
                'voices': voices,
            })
    except Exception as e:
        logging.getLogger(__name__).error('[TTS] List voices failed: %s', e, exc_info=True)
        return _error(str(e), 500)


@agent_matrix_bp.route('/tts/generate', methods=['POST'])
def tts_generate_edge():
    """Edge-TTS 文字转语音 — 直接返回 audio/mpeg 二进制流"""
    admin, err = _require_admin()
    if err: return err

    data = request.get_json()
    if not data or 'text' not in data:
        return _error(_('Text is required'))

    text = data.get('text')
    voice = data.get('voice', 'zh-CN-XiaoxiaoNeural')
    rate = data.get('rate', '+0%')

    try:
        from services.tts_service import text_to_speech_bytes
        import asyncio
        audio_bytes = asyncio.run(text_to_speech_bytes(text, voice, rate))
        from flask import Response
        return Response(
            audio_bytes,
            mimetype='audio/mpeg',
            headers={
                'Content-Disposition': 'attachment; filename="tts_output.mp3"',
                'X-TTS-Voice': voice,
                'X-TTS-Chars': str(len(text)),
            }
        )
    except Exception as e:
        return _error(str(e), 500)


# ============================================================
# 5.5 Discussion Mode — Multi-Agent Collaborative Chat
# ============================================================

@agent_matrix_bp.route('/chat/discuss', methods=['POST'])
def chat_discuss_sse():
    """
    SSE streaming discussion endpoint — multi-agent collaborative orchestration.

    Request: { "message": "...", "session_id": "(optional)" }
    Response: text/event-stream with event types:
      - phase: discussion phase indicator (planning/review/revision/decision/execution)
      - message: agent output content
      - needs_approval: JSON parse failed, requires manual intervention
      - warning: degradation notice (e.g. single-agent fallback)
      - error: unrecoverable error
      - done: discussion complete

    Total timeout: 300s enforced at Flask level.
    """
    admin, err = _require_admin()
    if err:
        return err

    ai_err = _check_ai_access()
    if ai_err:
        return ai_err

    data = request.get_json(force=True) or {}
    message = (data.get('message') or '').strip()
    if not message:
        return _error('Message cannot be empty')

    session_id = data.get('session_id', '')
    if not session_id:
        session_id = _m().create_session()
    # P1-F09: session_id 白名单校验（防路径穿越）
    elif not _m().is_valid_session_id(session_id):
        return _error(_('Invalid session ID'), 400)

    from agent_matrix.orchestrator import AgentOrchestrator
    from flask import Response, stream_with_context

    orchestrator = AgentOrchestrator(models_module=_m())

    def generate():
        try:
            for event in orchestrator.discuss_and_execute(
                instruction=message,
                user_id=admin.get('user_id', 0),
                session_id=session_id
            ):
                event['session_id'] = session_id
                yield f"data: {json.dumps(event)}\n\n"
        except Exception as e:
            import traceback
            traceback.print_exc()
            yield f"data: {json.dumps({'type': 'error', 'content': str(e), 'session_id': session_id})}\n\n"
        finally:
            yield "data: [DONE]\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
        }
    )


@agent_matrix_bp.route('/discuss/approve', methods=['POST'])
def discuss_approve():
    """
    Manual approval endpoint — triggered when Agent C outputs invalid JSON.

    Request: { "steps": [{"type":"...", "params":{...}}], "workflow_id": (optional) }
    Response: { "success": true, "message": "..." } or { "error": "..." }
    """
    admin, err = _require_admin()
    if err:
        return err

    ai_err = _check_ai_access()
    if ai_err:
        return ai_err

    data = request.get_json(force=True) or {}
    steps = data.get('steps')
    if not steps or not isinstance(steps, list):
        return _error('Steps must be a non-empty JSON array')

    workflow_id = data.get('workflow_id')

    exec_plan = {
        'approved': True,
        'confidence': 1.0,
        'reason': 'Manually approved by admin',
        'steps': steps,
    }
    if workflow_id:
        exec_plan['workflow_id'] = workflow_id

    from agent_matrix.orchestrator import AgentOrchestrator
    orchestrator = AgentOrchestrator(models_module=_m())

    try:
        result = orchestrator._trigger_dag_from_plan(exec_plan, admin.get('user_id', 0))
        return _success({'message': result})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return _error(f'Execution failed: {e}', 500)


# ============================================================
# 6. 初始化
# ============================================================

def init_agent_matrix(app):
    """初始化 Agent 矩阵系统（由 admin/app.py 调用）"""
    _m().init_agent_matrix_tables()
    _m().seed_default_agents()
    try:
        from agent_matrix.seed_prompts import seed_prompts
        seed_prompts()
    except Exception as e:
        print(f'[Agent Matrix] ⚠️ seed_prompts failed: {e}')
    app.register_blueprint(agent_matrix_bp)
    print(_('[Agent Matrix] ✅ Database + seed data has been initialized'))
    print(f'[Agent Matrix] 📋 API: /admin/agent-matrix/*')
    return _m()
