#!/usr/bin/env python3
"""project_workspace/services/doc_processor.py — Document processing pipeline.

Handles file validation, text extraction, chunking, embedding generation,
and citation extraction. All operations are scoped to a project_id.

可选依赖（PyMuPDF / python-docx / python-pptx）缺失时给出友好提示，
提取阶段返回带 [WARNING] 的原文而非空串，保证上传流程不中断。
"""

import hashlib
import logging
import os
import re

logger = logging.getLogger('project_workspace.doc_processor')

# 可选依赖探测：缺库时记录日志（不影响插件加载）
for _lib, _name in (('fitz', 'PyMuPDF'), ('docx', 'python-docx'), ('pptx', 'python-pptx')):
    try:
        __import__(_lib)
    except ImportError:
        logger.warning('Optional dependency "%s" is not installed (run: pip install %s)', _name, _name)

ALLOWED_EXTENSIONS = {
    'pdf', 'docx', 'txt', 'md', 'pptx', 'xlsx', 'csv',
}

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))


def resolve_storage_dir(config: dict) -> str:
    """插件文件存储目录：系统根目录 + config.storage_dir（默认 data/project_workspace/）。"""
    rel = (config or {}).get('storage_dir', 'data/project_workspace/')
    return os.path.join(PROJECT_ROOT, rel)


class DocProcessor:
    """Process uploaded documents: extract text, chunk, embed, and store."""

    def __init__(self, config: dict):
        self._config = config or {}
        self._embedder = None

    @property
    def _embed(self):
        if self._embedder is None:
            from .embedding import EmbeddingService
            self._embedder = EmbeddingService(self._config)
        return self._embedder

    @property
    def _chunk_size(self) -> int:
        return int(self._config.get('chunk_size', 1000))

    @property
    def _chunk_overlap(self) -> int:
        return int(self._config.get('chunk_overlap', 200))

    def validate_file(self, filename: str, file_size: int) -> tuple:
        """Validate file extension and size. Returns (ok, error_msg)."""
        ext = os.path.splitext(filename)[1].lstrip('.').lower()
        if ext not in ALLOWED_EXTENSIONS:
            return False, 'Unsupported file type: .%s' % ext
        max_size = int(self._config.get('max_file_size_mb', 50)) * 1024 * 1024
        if file_size > max_size:
            return False, 'File exceeds maximum size of %dMB' % (max_size // 1048576)
        return True, ''

    def extract_text(self, filepath: str, ext: str) -> str:
        """Extract text content from a file. 缺库时返回带 [WARNING] 的文本而非空串。"""
        ext = ext.lower()
        try:
            if ext in ('txt', 'md'):
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            elif ext == 'pdf':
                return self._extract_pdf(filepath)
            elif ext == 'docx':
                return self._extract_docx(filepath)
            elif ext == 'pptx':
                return self._extract_pptx(filepath)
            elif ext == 'csv':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            else:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        except Exception as e:
            logger.error('text extraction failed for %s: %s', filepath, e)
            return '[WARNING] text extraction failed: %s' % e

    def _extract_pdf(self, filepath: str) -> str:
        try:
            import fitz
        except ImportError:
            logger.warning('PyMuPDF not installed for PDF extraction')
            return '[WARNING] PDF extraction requires PyMuPDF (run: pip install PyMuPDF)'
        text = ''
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text()
        doc.close()
        return text

    def _extract_docx(self, filepath: str) -> str:
        try:
            from docx import Document
        except ImportError:
            logger.warning('python-docx not installed for DOCX extraction')
            return '[WARNING] DOCX extraction requires python-docx (run: pip install python-docx)'
        doc = Document(filepath)
        return '\n'.join(p.text for p in doc.paragraphs)

    def _extract_pptx(self, filepath: str) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning('python-pptx not installed for PPTX extraction')
            return '[WARNING] PPTX extraction requires python-pptx (run: pip install python-pptx)'
        prs = Presentation(filepath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    texts.append(shape.text)
        return '\n'.join(texts)

    # -- 块类型识别 ------------------------------------------------------

    def _looks_like_table_line(self, s: str) -> bool:
        return s.startswith('|') and s.endswith('|') and '|' in s[1:-1]

    def _segment_blocks(self, lines: list) -> list:
        """把原始行聚合为带类型的块：heading / table / code / paragraph。

        Returns: list[(block_type, content, section_title)]
        """
        blocks = []
        i, n = 0, len(lines)
        while i < n:
            line = lines[i]
            stripped = line.strip()

            # 代码块 ```...```：整体成块不切分
            if stripped.startswith('```'):
                buf = [line]
                i += 1
                while i < n and not lines[i].strip().startswith('```'):
                    buf.append(lines[i])
                    i += 1
                if i < n:
                    buf.append(lines[i])
                    i += 1
                blocks.append(('code', '\n'.join(buf), ''))
                continue

            # 标题块：≤3 级，作为 section_title 元数据
            hm = re.match(r'^(#{1,3})\s+(.*)$', stripped)
            if hm:
                blocks.append(('heading', stripped, stripped[:255]))
                i += 1
                continue

            # 表格块：连续 |...| 行整体成块不切分
            if self._looks_like_table_line(stripped):
                buf = [line]
                i += 1
                while i < n and self._looks_like_table_line(lines[i].strip()):
                    buf.append(lines[i])
                    i += 1
                blocks.append(('table', '\n'.join(buf), ''))
                continue

            # 普通段落：收集到空行或下一个块起点为止
            if stripped == '':
                i += 1
                continue
            buf = [line]
            i += 1
            while i < n:
                s = lines[i].strip()
                if (s == '' or s.startswith('```')
                        or re.match(r'^#{1,3}\s', s)
                        or self._looks_like_table_line(s)):
                    break
                buf.append(lines[i])
                i += 1
            blocks.append(('paragraph', '\n'.join(buf), ''))
        return blocks

    @staticmethod
    def _make_chunk(idx: int, content: str, section_title: str) -> dict:
        return {
            'chunk_index': idx,
            'content': content.strip(),
            'token_count': len(content.split()),
            'section_title': section_title,
        }

    def chunk_text(self, text: str) -> list:
        """按块类型切分文本，返回带 section_title 元数据的 chunks。

        - heading（#/##/###）：记录标题层级（≤3 级作为 section_title 元数据）
        - table / code：整体成块不切分
        - 普通段落：按 chunk_size 切分并保留 overlap 窗口
        """
        if not text.strip():
            return []
        blocks = self._segment_blocks(text.split('\n'))
        chunks = []
        current_idx = 0
        section_title = ''
        current_chunk = ''

        for btype, content, title in blocks:
            if title:
                section_title = title

            if btype in ('table', 'code'):
                # 整块独立成 chunk，不切分
                if current_chunk.strip():
                    chunks.append(self._make_chunk(current_idx, current_chunk, section_title))
                    current_idx += 1
                    current_chunk = ''
                chunks.append(self._make_chunk(current_idx, content, section_title))
                current_idx += 1
                continue

            # heading / paragraph：并入当前块（超长时切块并保留 overlap）
            if current_chunk.strip():
                if len(current_chunk) + len(content) > self._chunk_size:
                    chunks.append(self._make_chunk(current_idx, current_chunk, section_title))
                    current_idx += 1
                    overlap = current_chunk[-self._chunk_overlap:] if self._chunk_overlap > 0 else ''
                    current_chunk = (overlap + '\n' + content).strip()
                else:
                    current_chunk += '\n' + content
            else:
                current_chunk = content

        if current_chunk.strip():
            chunks.append(self._make_chunk(current_idx, current_chunk, section_title))

        return chunks

    def extract_citations(self, text: str) -> list:
        """Extract citation references from text using regex patterns."""
        citations = []
        seen = set()

        for m in re.finditer(r'\(([A-Z][a-zA-Z\-\']+(?:\s+et\s+al\.?)?),\s*(\d{4})\)', text):
            key = '%s_%s' % (m.group(1), m.group(2))
            if key not in seen:
                seen.add(key)
                citations.append({
                    'citation_key': key,
                    'authors': [m.group(1).strip()],
                    'year': int(m.group(2)),
                    'raw_text': m.group(0),
                })

        for m in re.finditer(r'\[(\d+(?:[,\-\s]\d+)*)\]', text):
            key = 'ref_%s' % m.group(1)
            if key not in seen:
                seen.add(key)
                citations.append({'citation_key': key, 'raw_text': m.group(0)})

        for m in re.finditer(r'(10\.\d{4,}/[^\s\]]+)', text):
            doi = m.group(1).rstrip('.')
            key = 'doi_%s' % hashlib.md5(doi.encode()).hexdigest()[:8]
            if key not in seen:
                seen.add(key)
                citations.append({'citation_key': key, 'doi': doi, 'raw_text': doi})

        return citations

    def process_document(self, doc_id: str, project_id: str, filepath: str,
                         filename: str, conn) -> bool:
        """Full processing pipeline: extract -> chunk -> embed -> store.

        本方法只 execute 不 commit，由调用方（_process_document_task）统一提交，
        保证 documents.status 与 projects.doc_count 在同一事务内。
        """
        ext = os.path.splitext(filename)[1].lstrip('.').lower()
        try:
            text = self.extract_text(filepath, ext)
            if not text:
                conn.execute(
                    "UPDATE documents SET status = 'failed', error_msg = 'No text could be extracted'"
                    " WHERE id = ?", (doc_id,)
                )
                return False

            word_count = len(text.split())
            char_count = len(text)
            conn.execute(
                "UPDATE documents SET word_count = ?, char_count = ?, status = 'processing'"
                " WHERE id = ?", (word_count, char_count, doc_id)
            )

            chunks = self.chunk_text(text)
            if not chunks:
                conn.execute(
                    "UPDATE documents SET status = 'failed', error_msg = 'Chunking produced empty result'"
                    " WHERE id = ?", (doc_id,)
                )
                return False

            for chunk in chunks:
                content = chunk['content']
                vec = self._embed.embed(content) if self._embed.is_ready() else None
                if vec:
                    vec_literal = '[' + ','.join(repr(v) for v in vec) + ']'
                    conn.execute(
                        "INSERT INTO document_chunks"
                        " (document_id, project_id, chunk_index, content, embedding,"
                        "  token_count, section_title)"
                        " VALUES (?, ?, ?, ?, ?::vector, ?, ?)",
                        (doc_id, project_id, chunk['chunk_index'], content,
                         vec_literal, chunk['token_count'], chunk['section_title'])
                    )
                else:
                    conn.execute(
                        "INSERT INTO document_chunks"
                        " (document_id, project_id, chunk_index, content, token_count, section_title)"
                        " VALUES (?, ?, ?, ?, ?, ?)",
                        (doc_id, project_id, chunk['chunk_index'], content,
                         chunk['token_count'], chunk['section_title'])
                    )

            citations = self.extract_citations(text)
            for cit in citations:
                conn.execute(
                    "INSERT INTO citations"
                    " (document_id, project_id, citation_key, authors, year, doi, raw_text)"
                    " VALUES (?, ?, ?, ?, ?, ?, ?)"
                    " ON CONFLICT (document_id, citation_key) DO NOTHING",
                    (doc_id, project_id, cit['citation_key'],
                     cit.get('authors', []), cit.get('year'),
                     cit.get('doi'), cit.get('raw_text'))
                )

            conn.execute(
                "UPDATE documents SET status = 'ready', processed_at = now(),"
                " page_count = ?, updated_at = now()"
                " WHERE id = ?",
                (len(chunks), doc_id)
            )
            return True

        except Exception as e:
            logger.error('document processing failed for %s: %s', doc_id, e)
            try:
                conn.execute(
                    "UPDATE documents SET status = 'failed', error_msg = ?"
                    " WHERE id = ?", (str(e)[:500], doc_id)
                )
            except Exception:
                pass
            return False

    def cleanup_stale_documents(self):
        """Archive documents older than retention_days (daily job)."""
        retention = int(self._config.get('retention_days', 730))
        from ..models import get_db
        conn = get_db()
        try:
            conn.execute(
                "UPDATE documents SET status = 'archived'"
                " WHERE status = 'ready'"
                " AND uploaded_at < CURRENT_DATE - INTERVAL '1 day' * ?",
                (retention,)
            )
            affected = conn.execute(
                "SELECT COUNT(*) AS n FROM documents WHERE status = 'archived'"
            ).fetchone()
            conn.commit()
            logger.info('cleanup archived %d stale documents', affected['n'] or 0)
        except Exception as e:
            logger.error('cleanup failed: %s', e)
            conn.rollback()
        finally:
            conn.close()
